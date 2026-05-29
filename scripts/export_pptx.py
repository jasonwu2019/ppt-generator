#!/usr/bin/env python3
"""
PPTX Export Script - Convert HTML presentation slides to PPTX format.

Usage:
    python export_pptx.py <input_html> <output_pptx>

Dependencies:
    python-pptx, Pillow

The script parses <section class="slide"> elements from the HTML file,
extracts text content from each slide, and creates a PPTX file with
16:9 slide dimensions and basic text formatting.

Note: This creates text-based PPTX slides. Complex layouts, images,
and CSS effects cannot be directly translated to PPTX format.
For best visual fidelity, use the HTML version for presentation.
"""

import sys
import os
import re
import html as html_lib
from html.parser import HTMLParser
from pathlib import Path


class SlideExtractor(HTMLParser):
    """Extract slide content from HTML presentation."""

    def __init__(self):
        super().__init__()
        self.slides = []
        self.current_slide = None
        self.in_slide = False
        self.in_style = False
        self.in_script = False
        self.text_buffer = []
        self.current_tag = None
        self.skip_depth = 0

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)

        if tag in ('style', 'script'):
            if tag == 'style':
                self.in_style = True
            else:
                self.in_script = True
            return

        if self.in_style or self.in_script:
            self.skip_depth += 1
            return

        # Detect slide sections
        if tag == 'section' and 'slide' in attrs_dict.get('class', ''):
            self.current_slide = {'type': 'content', 'title': '', 'items': []}
            self.in_slide = True
            slide_type = attrs_dict.get('data-slide-type', '')
            if slide_type:
                self.current_slide['type'] = slide_type

        # Track heading tags
        if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
            self.current_tag = tag

        # Track list/paragraph items
        if tag == 'p':
            self.current_tag = 'p'

    def handle_endtag(self, tag):
        if tag in ('style', 'script'):
            if tag == 'style':
                self.in_style = False
            else:
                self.in_script = False
            return

        if self.in_style or self.in_script:
            self.skip_depth -= 1
            return

        if tag == 'section' and self.in_slide:
            if self.current_slide:
                self.slides.append(self.current_slide)
            self.current_slide = None
            self.in_slide = False

        if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p'):
            if self.text_buffer and self.current_slide:
                text = ' '.join(self.text_buffer).strip()
                if text:
                    if tag in ('h1', 'h2') and not self.current_slide['title']:
                        self.current_slide['title'] = text
                    else:
                        self.current_slide['items'].append(text)
            self.text_buffer = []
            self.current_tag = None

    def handle_data(self, data):
        if self.in_style or self.in_script:
            return
        if self.in_slide and self.current_tag:
            cleaned = data.strip()
            if cleaned:
                self.text_buffer.append(cleaned)


def extract_slides(html_content):
    """Extract slide content from HTML string."""
    parser = SlideExtractor()
    parser.feed(html_content)
    return parser.slides


def create_pptx(slides, output_path):
    """Create PPTX file from extracted slides."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        print("ERROR: python-pptx is not installed.")
        print("Install it with: pip install python-pptx")
        sys.exit(1)

    prs = Presentation()

    # Set 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color definitions
    BLUE = RGBColor(0x00, 0x52, 0xD9)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x13, 0x1B, 0x2E)
    SLATE = RGBColor(0x47, 0x53, 0x69)

    # Mapping of slide types to backgrounds
    TYPE_BG = {
        'cover': BLUE,
        'ending': RGBColor(0x00, 0x5E, 0xFF),
        'toc': RGBColor(0xFA, 0xF8, 0xFF),
        'content': WHITE,
    }

    slide_count = 0
    for slide_data in slides:
        slide_count += 1
        slide_type = slide_data.get('type', 'content')
        bg_color = TYPE_BG.get(slide_type, WHITE)
        is_dark_bg = slide_type in ('cover', 'ending')

        # Use blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Set slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Add title
        title_text = slide_data.get('title', '')
        if title_text:
            # Position: centered or left-aligned
            if slide_type in ('cover', 'ending'):
                # Large centered title for cover/ending
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(11.333)
                height = Inches(2.5)
                alignment = PP_ALIGN.LEFT
            else:
                # Standard header
                left = Inches(1)
                top = Inches(0.5)
                width = Inches(11.333)
                height = Inches(1)
                alignment = PP_ALIGN.LEFT

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.alignment = alignment

            if is_dark_bg:
                p.font.color.rgb = WHITE
                p.font.size = Pt(44 if slide_type == 'cover' else 36)
            else:
                p.font.color.rgb = BLUE
                p.font.size = Pt(32)

            p.font.bold = True
            if slide_type == 'ending':
                p.font.italic = True

        # Add content items
        items = slide_data.get('items', [])
        if items:
            if slide_type in ('toc',):
                # Two-column layout for TOC
                mid = (len(items) + 1) // 2
                col1 = items[:mid]
                col2 = items[mid:]

                for col_idx, col_items in enumerate((col1, col2)):
                    left = Inches(1 + col_idx * 6)
                    top = Inches(1.8)
                    width = Inches(5.5)
                    height = Inches(5)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    for i, item in enumerate(col_items):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()

                        num = i + 1 + (col_idx * mid)
                        p.text = f"  {num}.  {item}"
                        p.font.size = Pt(20)
                        p.font.color.rgb = DARK
                        p.space_after = Pt(16)
            else:
                # Regular content list
                left = Inches(1.5)
                top = Inches(2 if title_text else 1)
                width = Inches(10)
                height = Inches(5)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                for i, item in enumerate(items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    # Clean HTML entities
                    cleaned = html_lib.unescape(re.sub(r'<[^>]+>', '', item))
                    p.text = cleaned
                    p.font.size = Pt(18)
                    p.font.color.rgb = SLATE if not is_dark_bg else WHITE
                    p.space_after = Pt(12)

        # Add slide number for non-cover slides
        if slide_type not in ('cover', 'ending'):
            left = Inches(12)
            top = Inches(7)
            width = Inches(1)
            height = Inches(0.3)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_count)
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            p.alignment = PP_ALIGN.RIGHT

    # Save
    prs.save(output_path)
    return slide_count


def main():
    if len(sys.argv) < 3:
        print("Usage: python export_pptx.py <input_html> <output_pptx>")
        print("Example: python export_pptx.py presentation.html presentation.pptx")
        sys.exit(1)

    input_html = sys.argv[1]
    output_pptx = sys.argv[2]

    if not os.path.exists(input_html):
        print(f"ERROR: Input file not found: {input_html}")
        sys.exit(1)

    # Read HTML
    with open(input_html, 'r', encoding='utf-8') as f:
        html_content = f.read()

    # Extract slides
    slides = extract_slides(html_content)
    print(f"Found {len(slides)} slides in {input_html}")

    if not slides:
        print("WARNING: No slides found. Make sure the HTML uses <section class=\"slide\"> elements.")
        sys.exit(1)

    # Create PPTX
    count = create_pptx(slides, output_pptx)
    print(f"Created PPTX with {count} slides: {output_pptx}")
    print(f"Output size: {os.path.getsize(output_pptx) / 1024:.1f} KB")


if __name__ == '__main__':
    main()
